from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def parse_markdown_to_blocks(markdown: str) -> list[dict]:
    lines = [l.rstrip() for l in markdown.splitlines()]

    slides: list[dict] = []
    current: dict | None = None

    def flush() -> None:
        nonlocal current
        if current and (current["title"] or current["bullets"] or current["notes"]):
            slides.append(current)
        current = None

    for line in lines:
        if line.startswith("## "):
            flush()
            current = {"title": line[3:].strip(), "bullets": [], "notes": []}
            continue
        if current is None:
            continue

        if line.startswith("### "):
            current["bullets"].append(
                {"level": 0, "text": line[4:].strip(), "bold": True}
            )
            continue

        m = re.match(r"^\s*-\s+(.*)$", line)
        if m:
            b = m.group(1).strip()
            b = re.sub(r"\*\*(.*?)\*\*", r"\1", b)
            current["bullets"].append({"level": 0, "text": b, "bold": False})
            continue

        if line.strip() in {"", "---"}:
            continue

        current["notes"].append(line)

    flush()
    return slides


def first_h1(markdown: str, default: str = "Présentation") -> str:
    for line in markdown.splitlines():
        if line.startswith("# "):
            return line[2:].strip()
    return default


def build_pptx(md_path: Path, out_path: Path) -> None:
    md = md_path.read_text(encoding="utf-8")
    blocks = parse_markdown_to_blocks(md)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Palette (inspirée du style SERAP-UAC)
    BG = RGBColor(10, 12, 16)  # #0a0c10
    CARD = RGBColor(17, 24, 39)  # #111827
    ACCENT = RGBColor(139, 92, 246)  # #8b5cf6
    ACCENT_2 = RGBColor(6, 182, 212)  # #06b6d4
    TEXT = RGBColor(243, 244, 246)  # #f3f4f6
    MUTED = RGBColor(156, 163, 175)  # #9ca3af

    def apply_background(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def add_accent_bar(slide, title_text: str) -> None:
        # Top bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.85)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Thin secondary line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.85), prs.slide_width, Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_2
        line.line.fill.background()

        # Title text box on top bar (instead of placeholder title)
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), prs.slide_width - Inches(1.2), Inches(0.6))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(34)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

    def style_body_textframe(tf) -> None:
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.color.rgb = TEXT
            if p.text and p.level == 0:
                p.font.size = Pt(22)
            elif p.text:
                p.font.size = Pt(18)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    apply_background(slide)
    # Hide default placeholders by moving them off-canvas (python-pptx can't "delete" reliably)
    slide.shapes.title.left = prs.slide_width
    slide.shapes.title.top = prs.slide_height
    slide.placeholders[1].left = prs.slide_width
    slide.placeholders[1].top = prs.slide_height

    add_accent_bar(slide, first_h1(md, default="Présentation SNMP"))

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), prs.slide_width - Inches(1.6), Inches(1.4))
    stf = subtitle_box.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    p.text = "SNMP (Simple Network Management Protocol)"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = TEXT
    p2 = stf.add_paragraph()
    p2.text = "Synthèse structurée + comparatif de versions + points sécurité"
    p2.runs[0].font.size = Pt(18)
    p2.runs[0].font.color.rgb = MUTED

    # Content slides
    bullet_layout = prs.slide_layouts[1]  # Title and Content
    for block in blocks:
        slide = prs.slides.add_slide(bullet_layout)
        apply_background(slide)

        # Hide default title placeholder and replace with accent bar
        slide.shapes.title.left = prs.slide_width
        slide.shapes.title.top = prs.slide_height
        add_accent_bar(slide, block["title"])

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        bullets = block["bullets"]
        if not bullets and block["notes"]:
            bullets = [{"level": 0, "text": " ".join(block["notes"]).strip(), "bold": False}]

        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b["text"]
            p.level = int(b.get("level", 0))
            p.font.size = Pt(22 if p.level == 0 else 18)
            p.font.color.rgb = TEXT

            if b.get("bold"):
                for run in p.runs:
                    run.font.bold = True

        # Content container "card" effect behind text
        # (add after text; send to back among shapes but above background)
        try:
            ph = slide.shapes.placeholders[1]
            left, top, width, height = ph.left, ph.top, ph.width, ph.height
            pad = Inches(0.1)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - pad, top - pad, width + 2 * pad, height + 2 * pad)
            card.fill.solid()
            card.fill.fore_color.rgb = CARD
            card.fill.transparency = 0.08
            card.line.color.rgb = ACCENT
            card.line.width = Pt(1.25)
            # Move card behind placeholder content by reordering the spTree
            sp = card._element
            spTree = slide.shapes._spTree
            spTree.remove(sp)
            spTree.insert(2, sp)  # after background & bar
        except Exception:
            pass

        if block["notes"]:
            notes = slide.notes_slide.notes_text_frame
            notes.clear()
            notes.text = "\n".join(block["notes"]).strip()

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Génère un fichier .pptx à partir d'un markdown structuré (# titre, ## sections, - puces)."
    )
    p.add_argument(
        "--src",
        type=Path,
        default=Path("presentation_SNMP_nouvelle.md"),
        help="Chemin du fichier markdown source.",
    )
    p.add_argument(
        "--out",
        type=Path,
        default=Path("presentation_SNMP_nouvelle.pptx"),
        help="Chemin du fichier .pptx de sortie.",
    )
    return p.parse_args()


if __name__ == "__main__":
    args = parse_args()
    build_pptx(args.src, args.out)
    print(f"WROTE {args.out.resolve()}")
